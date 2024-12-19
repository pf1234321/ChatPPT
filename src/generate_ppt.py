from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# 创建一个新的Presentation对象
prs = Presentation()

# 幻灯片1：标题和文本
slide1_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide1_layout)
title1 = slide1.shapes.title
title1.text = "示例演示文稿"
subtitle1 = slide1.placeholders[1]
subtitle1.text = "包含文本、图片、表格和图表"

# 幻灯片2：文本和图片
slide2_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(slide2_layout)
title2 = slide2.shapes.title
title2.text = "文本与图片"
left = Inches(1)
top = Inches(1)
width = Inches(3)
height = Inches(1)
txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "这是一些示例文本，展示如何在幻灯片中添加文本和图片。"
p.font.size = Pt(14)
# 添加图片（请确保有一张名为'example.jpg'的图片在当前目录下）
pic = slide2.shapes.add_picture('example.jpg', Inches(4), Inches(1), width=Inches(3), height=Inches(3))

# 幻灯片3：表格
slide3_layout = prs.slide_layouts[1]
slide3 = prs.slides.add_slide(slide3_layout)
title3 = slide3.shapes.title
title3.text = "表格"
# 创建表格数据
table_data = [
    ["姓名", "年龄", "性别"],
    ["张三", 25, "男"],
    ["李四", 30, "女"],
    ["王五", 28, "男"]
]
# 创建表格形状
left = Inches(1)
top = Inches(1)
width = Inches(5)
height = Inches(2)
table = slide3.shapes.add_table(rows=len(table_data), cols=len(table_data[0]), left=left, top=top, width=width, height=height).table
# 填充表格数据
for row in range(len(table_data)):
    for col in range(len(table_data[row])):
        cell = table.cell(row, col)
        cell.text = str(table_data[row][col])

# 幻灯片4：图表
slide4_layout = prs.slide_layouts[1]
slide4 = prs.slides.add_slide(slide4_layout)
title4 = slide4.shapes.title
title4.text = "图表"
# 创建图表数据
chart_data = ChartData()
chart_data.categories = ['苹果', '香蕉', '橙子']
chart_data.add_series('销量', (10, 20, 15))
# 创建图表形状
left = Inches(1)
top = Inches(1)
width = Inches(5)
height = Inches(3)
chart = slide4.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
).chart

# 保存文件
prs.save("example.pptx")
